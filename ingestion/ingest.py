#!/usr/bin/env python3
"""
Document ingestion for the AI stack document brain.
Reads documents from /documents/<acl_folder>/ (RECURSIVE — subfolders included),
chunks, embeds, and stores them in Qdrant with ACL tags. Runs on a schedule inside
the ingestion container (every INGEST_INTERVAL_SECONDS, default 900s) and is also
runnable manually:  docker exec ingestion python3 /app/ingest.py

File types: pdf, docx, txt, md, markdown, rtf, html/htm, csv, xlsx, pptx, odt.
Legacy binary .doc/.xls and anything else: skipped as unsupported. (OCR for
scanned PDFs belongs to the dedicated-pipeline upgrade — see roadmap.)

Scheduled-worker support:
- sha256 manifest at /app/.ingest-manifest.json — unchanged document set = free no-op cycle.
- Deletion reconciliation — files removed or changed since the last successful run
  lose all their Qdrant points before re-ingestion (no stale answers from deleted
  documents, no orphaned tail chunks from shortened files).
- The manifest is saved ONLY on a clean pass (0 failures) so failed cycles retry.
Idempotent: deterministic point IDs (uuid5 on relpath:chunk) make re-runs overwrite.
NOTE: point IDs and payload 'source' use the relpath (company/hr/policy.txt) —
same-named files in different subfolders do NOT collide.
"""
import os, sys, uuid, glob, json, hashlib
import requests

# --- Config (from environment) ---
QDRANT_URL = os.environ.get("QDRANT_URL", "http://qdrant:6333")
QDRANT_API_KEY = os.environ["QDRANT_API_KEY"]
EMBED_URL = os.environ.get("EMBED_URL", "http://embeddings:80")
COLLECTION = "company_docs"
DOCS_ROOT = "/documents"
# Folder name directly under DOCS_ROOT == the ACL tag applied (subfolders included).
ACL_FOLDERS = ("company", "executive")
CHUNK_SIZE = 512        # characters per chunk (simple char chunking)
CHUNK_OVERLAP = 64
VECTOR_SIZE = 1024      # bge-m3 output dimension — must match embeddings service

HEADERS = {"api-key": QDRANT_API_KEY, "Content-Type": "application/json"}

MANIFEST_PATH = "/app/.ingest-manifest.json"   # lives on the host mount; survives recreates

# Outcome tracking for the end-of-run summary.
STATS = {"ok": 0, "skipped": 0, "failed": 0, "chunks": 0, "stale": 0}
FAILURES = []  # (path, reason) for the failure summary


def extract_text(path):
    """Pull raw text from a document. Returns a string, or None to skip.

    Returns None (skip) only for genuinely unsupported types; parse errors raise
    and are counted by main()'s per-file handler.
    """
    ext = path.lower().rsplit(".", 1)[-1]

    # --- PDF ---
    if ext == "pdf":
        from pypdf import PdfReader
        reader = PdfReader(path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    # --- Modern Word ---
    if ext == "docx":
        import docx
        doc = docx.Document(path)
        return "\n".join(p.text for p in doc.paragraphs)

    # --- Plain text family (direct read, tolerate weird encodings) ---
    if ext in ("txt", "md", "markdown", "csv"):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    # --- HTML / HTM: strip tags, keep body text (don't embed the markup) ---
    if ext in ("html", "htm"):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            raw = f.read()
        return _html_to_text(raw)

    # --- RTF: strip control words to approximate plain text ---
    if ext == "rtf":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return _rtf_to_text(f.read())

    # --- Excel (xlsx via openpyxl) ---
    if ext == "xlsx":
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        lines = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines)

    # --- PowerPoint (pptx via python-pptx) ---
    if ext == "pptx":
        from pptx import Presentation
        prs = Presentation(path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
        return "\n".join(lines)

    # --- OpenDocument text (odt) ---
    if ext == "odt":
        return _odt_to_text(path)

    # --- Genuinely unsupported here: legacy binary .doc/.xls, and anything else ---
    print(f"  SKIP (unsupported type): {path}")
    return None


def _html_to_text(raw):
    """HTML -> visible text. BeautifulSoup if available, else a regex fallback.

    WHY: we want the page's body text embedded, not the <div>/<script> markup.
    """
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
        lines = [ln.strip() for ln in text.splitlines()]
        return "\n".join(ln for ln in lines if ln)
    except ImportError:
        import re
        raw = re.sub(r"(?s)<(script|style).*?</\1>", " ", raw)
        raw = re.sub(r"(?s)<[^>]+>", " ", raw)
        import html as _html
        return _html.unescape(re.sub(r"\s+", " ", raw)).strip()


def _rtf_to_text(raw):
    """RTF -> approximate plain text by stripping control words/groups.

    Good enough for simple RTF; complex layouts may leave residue.
    """
    import re
    raw = re.sub(r"\{\\[^{}]*\}", " ", raw)          # font/color tables and other {\...} groups
    raw = re.sub(r"\\[a-z]+-?\d* ?", " ", raw)        # control words like \par \b \fs24
    raw = raw.replace("{", " ").replace("}", " ").replace("\\", " ")
    return re.sub(r"\s+", " ", raw).strip()


def _odt_to_text(path):
    """ODT -> text. An .odt is a zip; the content is content.xml. Strip its tags."""
    import zipfile, re
    import xml.etree.ElementTree as ET
    with zipfile.ZipFile(path) as z:
        xml_data = z.read("content.xml")
    try:
        root = ET.fromstring(xml_data)
        text = " ".join(t.strip() for t in root.itertext() if t.strip())
        return re.sub(r"\s+", " ", text).strip()
    except ET.ParseError:
        return re.sub(r"(?s)<[^>]+>", " ", xml_data.decode("utf-8", "ignore")).strip()


def chunk_text(text):
    """Simple overlapping character chunker. Not fancy, works fine for business docs."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + CHUNK_SIZE
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - CHUNK_OVERLAP
    return chunks


def embed(texts):
    """Embed a list of strings via TEI, batched, with one retry per batch.

    WHY batched: one giant POST trips 413 Payload Too Large.
    WHY retry: a transient error shouldn't fail a whole file on a big corpus.
    """
    import time
    all_vectors = []
    BATCH = 32
    for i in range(0, len(texts), BATCH):
        part = texts[i:i+BATCH]
        for attempt in (1, 2):  # one retry
            try:
                resp = requests.post(f"{EMBED_URL}/embed", json={"inputs": part}, timeout=180)
                resp.raise_for_status()
                all_vectors.extend(resp.json())
                break
            except requests.RequestException:
                if attempt == 2:
                    raise
                time.sleep(2)
    return all_vectors


def ensure_collection():
    """Create the collection if it doesn't exist."""
    r = requests.get(f"{QDRANT_URL}/collections/{COLLECTION}", headers=HEADERS, timeout=30)
    if r.status_code == 200:
        return
    payload = {"vectors": {"size": VECTOR_SIZE, "distance": "Cosine"}}
    r = requests.put(f"{QDRANT_URL}/collections/{COLLECTION}", headers=HEADERS, json=payload, timeout=30)
    r.raise_for_status()
    print(f"Created collection '{COLLECTION}' (size={VECTOR_SIZE}, cosine)")


def doc_point_id(source_rel, chunk_index):
    """Deterministic UUID from relpath+chunk so re-runs overwrite instead of duplicate,
    and same-named files in different subfolders don't collide."""
    return str(uuid.uuid5(uuid.NAMESPACE_DNS, f"{source_rel}:{chunk_index}"))


def ingest_file(path, acl):
    source_rel = os.path.relpath(path, DOCS_ROOT)
    text = extract_text(path)
    if text is None:
        STATS["skipped"] += 1
        return 0
    if not text.strip():
        print(f"  SKIP (no text extracted): {source_rel}")
        STATS["skipped"] += 1
        return 0
    chunks = chunk_text(text)
    if not chunks:
        print(f"  SKIP (no chunks): {source_rel}")
        STATS["skipped"] += 1
        return 0
    print(f"Ingesting [{acl}] {source_rel} ... ({len(chunks)} chunks)")
    vectors = embed(chunks)
    points = []
    for i, (chunk, vec) in enumerate(zip(chunks, vectors)):
        points.append({
            "id": doc_point_id(source_rel, i),
            "vector": vec,
            "payload": {
                "source": source_rel,
                "chunk_index": i,
                "acl": acl,
                "text": chunk,
            },
        })
    # Upsert in batches of 64
    for i in range(0, len(points), 64):
        batch = points[i:i+64]
        r = requests.put(f"{QDRANT_URL}/collections/{COLLECTION}/points?wait=true",
                         headers=HEADERS, json={"points": batch}, timeout=120)
        r.raise_for_status()
    STATS["ok"] += 1
    STATS["chunks"] += len(points)
    return len(points)


# ---------------- scheduled-worker support ----------------

def sha256_file(path):
    h = hashlib.new("sha256")
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(1 << 20), b""):
            h.update(block)
    return h.hexdigest()

def scan_documents():
    """{relpath: sha256} for every file under the ACL folders (same recursive glob
    ingestion uses, so the manifest and the ingestion loop always agree)."""
    manifest = {}
    for acl in ACL_FOLDERS:
        folder = os.path.join(DOCS_ROOT, acl)
        for path in sorted(glob.glob(os.path.join(folder, "**", "*.*"), recursive=True)):
            try:
                manifest[os.path.relpath(path, DOCS_ROOT)] = sha256_file(path)
            except OSError as e:
                print(f"  WARNING: cannot read {path}: {e}", file=sys.stderr)
    return manifest

def load_manifest():
    try:
        with open(MANIFEST_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    except (OSError, ValueError):
        return {}

def save_manifest(manifest):
    with open(MANIFEST_PATH, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=1, sort_keys=True)

def delete_by_source(source_rel):
    """Remove all Qdrant points whose payload 'source' is <source_rel>."""
    body = {"filter": {"must": [{"key": "source", "match": {"value": source_rel}}]}}
    r = requests.post(f"{QDRANT_URL}/collections/{COLLECTION}/points/delete",
                      headers=HEADERS, json=body, timeout=60)
    r.raise_for_status()

def main():
    new_manifest = scan_documents()
    old_manifest = load_manifest()
    if old_manifest == new_manifest:
        print("No document changes since last run; nothing to do.")
        return

    ensure_collection()

    # Deletion reconciliation: files removed from disk, or changed since the last
    # successful run, lose ALL their points before re-ingestion.
    stale = [rel for rel in old_manifest
             if rel not in new_manifest or old_manifest[rel] != new_manifest[rel]]
    for rel in stale:
        delete_by_source(rel)
        print(f"Removed stale chunks for: {rel}")
    STATS["stale"] = len(stale)

    for acl in ACL_FOLDERS:
        folder = os.path.join(DOCS_ROOT, acl)
        if not os.path.isdir(folder):
            print(f"NOTE: folder not present, skipping: {folder}")
            continue
        for path in sorted(glob.glob(os.path.join(folder, "**", "*.*"), recursive=True)):
            try:
                ingest_file(path, acl)
            except Exception as e:
                # A roadblock we WANT to see: record it, keep going.
                print(f"  ERROR on {path}: {e}", file=sys.stderr)
                STATS["failed"] += 1
                FAILURES.append((path, str(e)[:200]))

    # Save the manifest ONLY on a clean pass. A manifest saved after failures
    # (e.g. embeddings not up yet on first boot) poisons state: every later cycle
    # no-ops while Qdrant stays empty/stale. Unsaved = retried next cycle.
    if STATS["failed"]:
        print(f"\n{STATS['failed']} file(s) failed — manifest NOT saved; retrying next cycle.", file=sys.stderr)
    else:
        save_manifest(new_manifest)

    print("\n===== INGEST SUMMARY =====")
    print(f"  Files ingested OK : {STATS['ok']}")
    print(f"  Files skipped     : {STATS['skipped']} (unsupported type or no text)")
    print(f"  Files FAILED      : {STATS['failed']}")
    print(f"  Total chunks      : {STATS['chunks']}")
    print(f"  Stale sources removed: {STATS['stale']}")
    if FAILURES:
        print("\n  Failure detail (first 40):")
        for p, reason in FAILURES[:40]:
            print(f"    {p} -> {reason}")
    # Show collection stats
    r = requests.get(f"{QDRANT_URL}/collections/{COLLECTION}", headers=HEADERS, timeout=30)
    info = r.json().get("result", {})
    print(f"\nCollection '{COLLECTION}' points count: {info.get('points_count')}")
    sys.exit(1 if STATS["failed"] else 0)

if __name__ == "__main__":
    main()
